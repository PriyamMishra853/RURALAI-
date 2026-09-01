"""
Render each slide of a .pptx to PNG from real shape geometry, and report any
text that overflows its box.

No LibreOffice/PowerPoint here, so this reads position, size, fill, line and
per-run text out of the package and draws it with PIL using the same font
family PowerPoint will use. Wrapping is approximated the same way PowerPoint
does it (greedy, on the shape's inner width), so the overflow numbers are
meaningful even though the pixels are not pixel-identical.
"""
import sys, os, io
from pptx import Presentation
from pptx.util import Emu
from PIL import Image, ImageDraw, ImageFont

PX = 110          # px per inch
EMU_IN = 914400.0
F = "C:/Windows/Fonts/"
FONTS = {(0, 0): "calibri.ttf", (1, 0): "calibrib.ttf",
         (0, 1): "calibrii.ttf", (1, 1): "calibriz.ttf"}
_cache = {}


def fnt(sz, b=False, i=False):
    key = (round(sz, 1), b, i)
    if key not in _cache:
        _cache[key] = ImageFont.truetype(F + FONTS[(int(b), int(i))], max(6, int(sz * PX / 72)))
    return _cache[key]


def emu2px(v):
    return (v or 0) / EMU_IN * PX


def rgb(c, default=(20, 20, 20)):
    try:
        if c and c.type is not None and c.rgb is not None:
            return tuple(c.rgb)
    except Exception:
        pass
    return default


def shape_fill(sh):
    try:
        f = sh.fill
        if f.type is not None and "SOLID" in str(f.type):
            return tuple(f.fore_color.rgb)
    except Exception:
        pass
    return None


def shape_line(sh):
    try:
        if sh.line.fill.type is not None and "SOLID" in str(sh.line.fill.type):
            w = sh.line.width.pt if sh.line.width else 1
            return tuple(sh.line.color.rgb), max(1, int(w * PX / 72))
    except Exception:
        pass
    return None, 0


def wrap_runs(runs, maxw):
    """runs = [(text, font, colour)] -> list of lines, each a list of fragments."""
    lines, cur, curw = [], [], 0
    for txt, f, col in runs:
        for word in txt.replace("\n", " \n ").split(" "):
            if word == "":
                continue
            if word == "\n":
                lines.append(cur); cur, curw = [], 0; continue
            piece = word + " "
            w = f.getlength(piece)
            if curw + w > maxw and cur:
                lines.append(cur); cur, curw = [], 0
            cur.append((piece, f, col)); curw += w
    if cur:
        lines.append(cur)
    return lines


def render(path, outdir, prefix):
    prs = Presentation(path)
    W = int(prs.slide_width / EMU_IN * PX)
    H = int(prs.slide_height / EMU_IN * PX)
    problems = []
    made = []

    for si, slide in enumerate(prs.slides, 1):
        img = Image.new("RGB", (W, H), (255, 255, 255))
        d = ImageDraw.Draw(img)

        for sh in slide.shapes:
            try:
                x, y = emu2px(sh.left), emu2px(sh.top)
                w, h = emu2px(sh.width), emu2px(sh.height)
            except Exception:
                continue
            if w <= 0 or h <= 0:
                continue

            if sh.shape_type is not None and "PICTURE" in str(sh.shape_type):
                try:
                    im = Image.open(io.BytesIO(sh.image.blob)).convert("RGB")
                    img.paste(im.resize((max(1, int(w)), max(1, int(h))), Image.LANCZOS),
                              (int(x), int(y)))
                except Exception:
                    d.rectangle([x, y, x + w, y + h], outline=(200, 200, 200))
                continue

            fill = shape_fill(sh)
            lcol, lw = shape_line(sh)
            if fill or lcol:
                rad = min(14, h / 2, w / 2)
                d.rounded_rectangle([x, y, x + w, y + h], radius=rad,
                                    fill=fill, outline=lcol, width=lw or 0)

            if not sh.has_text_frame:
                continue
            tf = sh.text_frame
            ml = emu2px(tf.margin_left); mr = emu2px(tf.margin_right)
            mt = emu2px(tf.margin_top); mb = emu2px(tf.margin_bottom)
            iw = max(4, w - ml - mr)
            cy = y + mt
            total = 0

            for p in tf.paragraphs:
                runs = []
                for r in p.runs:
                    sz = r.font.size.pt if r.font.size else 12
                    b = bool(r.font.bold); it = bool(r.font.italic)
                    runs.append((r.text, fnt(sz, b, it), rgb(r.font.color)))
                if not runs:
                    cy += 8; total += 8; continue
                sb = p.space_before.pt * PX / 72 if p.space_before else 0
                sa = p.space_after.pt * PX / 72 if p.space_after else 0
                cy += sb; total += sb
                lines = wrap_runs(runs, iw)
                lh = max(f.size for _, f, _ in runs) * 1.22
                if p.line_spacing and isinstance(p.line_spacing, float):
                    lh *= p.line_spacing
                for line in lines:
                    lw_px = sum(f.getlength(t) for t, f, _ in line)
                    if p.alignment is not None and "CENTER" in str(p.alignment):
                        cx = x + ml + (iw - lw_px) / 2
                    elif p.alignment is not None and "RIGHT" in str(p.alignment):
                        cx = x + ml + iw - lw_px
                    else:
                        cx = x + ml
                    for t, f, col in line:
                        d.text((cx, cy), t, font=f, fill=col)
                        cx += f.getlength(t)
                    cy += lh; total += lh
                cy += sa; total += sa

            avail = h - mt - mb
            if total > avail + 3:
                problems.append(
                    f"  slide {si}: '{sh.name}' text {total:.0f}px in {avail:.0f}px box "
                    f"(+{total - avail:.0f}px)  «{tf.text[:58].strip()}»")

        out = os.path.join(outdir, f"{prefix}-{si}.png")
        img.save(out)
        made.append(out)
    return made, problems


if __name__ == "__main__":
    src = sys.argv[1]
    pref = sys.argv[2] if len(sys.argv) > 2 else "qa"
    here = os.path.dirname(os.path.abspath(src)) or "."
    files, probs = render(src, here, pref)
    for f in files:
        print(f)
    print("\n--- OVERFLOW REPORT ---")
    print("\n".join(probs) if probs else "  none")
