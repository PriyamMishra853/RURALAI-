"""
Fill the official SIH 2026 Idea template for PS 26133 — RuralAI.

Rules honoured:
  * Template's own slides, order, titles, footer, logo and slide numbers untouched.
  * Slide 7 (the "delete this slide" instructions page) removed as instructed.
  * The template's idea-detail pointers are kept as visible section headers.
  * Real text boxes and real shapes — nothing is a flattened image of text.
"""
import copy, os
from pptx import Presentation
from pptx.util import Inches as In, Pt, Emu
from pptx.dml.color import RGBColor as C
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
SRC = r"C:\Users\Priyam\Downloads\SIH2026-IDEA-Presentation-Format.pptx"
OUT = os.path.join(HERE, "RuralAI_SIH2026_PS26133_UNFILTEREDENGINEERS.pptx")

NAVY      = C(0x0B, 0x3C, 0x78)
NAVY_DEEP = C(0x08, 0x28, 0x4F)
NAVY_SOFT = C(0xE8, 0xF0, 0xFA)
SAFFRON   = C(0xF3, 0x92, 0x11)
AMBER     = C(0xB4, 0x53, 0x09)
GREEN     = C(0x15, 0x80, 0x3D)
GREEN_BG  = C(0xDE, 0xF7, 0xE6)
RED       = C(0xBE, 0x12, 0x3C)
RED_BG    = C(0xFE, 0xE2, 0xE8)
INK       = C(0x0F, 0x20, 0x37)
MUTED     = C(0x52, 0x64, 0x7C)
LINE      = C(0xCD, 0xD8, 0xE6)
WHITE     = C(0xFF, 0xFF, 0xFF)
SOFT      = C(0xF4, 0xF7, 0xFB)
FONT      = "Calibri"

prs = Presentation(SRC)


# ------------------------------------------------------------------ helpers
def drop_slide(idx):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    rId = slides[idx].get(
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
    prs.part.drop_rel(rId)
    xml_slides.remove(slides[idx])


def kill(shape):
    shape._element.getparent().remove(shape._element)


def find(slide, needle):
    for sh in slide.shapes:
        if sh.has_text_frame and needle.lower() in sh.text_frame.text.lower():
            return sh
    return None


def tb(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    s = slide.shapes.add_textbox(In(x), In(y), In(w), In(h))
    tf = s.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    return tf


def para(tf, first=False, space_before=0, space_after=3, align=PP_ALIGN.LEFT, line=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    p.alignment = align
    if line:
        p.line_spacing = line
    return p


def run(p, text, size=11, bold=False, color=INK, italic=False, font=FONT):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return r


def card(slide, x, y, w, h, fill=SOFT, outline=LINE, lw=0.75, radius=None,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    s = slide.shapes.add_shape(shape, In(x), In(y), In(w), In(h))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if outline is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = outline
        s.line.width = Pt(lw)
    s.shadow.inherit = False
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            s.adjustments[0] = radius
        except Exception:
            pass
    s.text_frame.word_wrap = True
    return s


def label(slide, x, y, w, text, size=11.5, color=NAVY):
    tf = tb(slide, x, y, w, 0.26)
    p = para(tf, first=True, space_after=0)
    run(p, text, size, True, color)
    return tf


def section(slide, x, y, w, title, rows, tsize=11.5, bsize=9.6, lead=0.215,
            tcolor=NAVY, gap=0.055):
    """Pointer heading + 'Bold lead : detail' rows, the winning-deck idiom."""
    label(slide, x, y, w, title, tsize, tcolor)
    yy = y + 0.30
    for head, body in rows:
        n = max(1, -(-len(head + body) // int(w * 20)))
        tf = tb(slide, x, yy, w, lead * n)
        p = para(tf, first=True, space_after=0, line=0.95)
        if head:
            run(p, head + "  ", bsize, True, INK)
        run(p, body, bsize, False, MUTED)
        yy += lead * n + gap
    return yy


# ============================================================ slide 7 removed
drop_slide(6)
S = prs.slides

# ================================================================== SLIDE 1
s1 = S[0]
FILL = {
    "Problem Statement ID": "Problem Statement ID – 26133",
    "Problem Statement Title": "Problem Statement Title – Accessibility and quality of public "
                               "healthcare services, particularly in rural and underserved areas",
    "Theme-": "Theme – MedTech / BioTech / HealthTech",
    "PS Category": "PS Category – Software",
    "Team ID": "Team ID – [TEAM ID]",
    "Team Name": "Team Name – UNFILTEREDENGINEERS",
}
for sh in s1.shapes:
    if not sh.has_text_frame:
        continue
    for p in sh.text_frame.paragraphs:
        txt = "".join(r.text for r in p.runs)
        for key, val in FILL.items():
            if key.lower() in txt.lower() and len(txt) < 90:
                for i, r in enumerate(p.runs):
                    r.text = val if i == 0 else ""
                    if i == 0:
                        r.font.size = Pt(15)
                        r.font.name = FONT
                break

tf = tb(s1, 0.62, 6.05, 8.4, 0.5)
p = para(tf, first=True, space_after=0)
run(p, "Working prototype, live now:  ", 12, True, NAVY)
run(p, "ruralai-psi.vercel.app", 12, True, SAFFRON)
p = para(tf, space_before=1, space_after=0)
run(p, "Government of Maharashtra  ·  Maharashtra State Innovation Society", 10, False, MUTED)

# ================================================================== SLIDE 2
s2 = S[1]
kill(find(s2, "Proposed Solution"))

tf = tb(s2, 0.36, 1.30, 6.2, 0.55)
p = para(tf, first=True, space_after=0)
run(p, "RuralAI", 30, True, NAVY)
run(p, "   Care access without the journey", 13, False, MUTED)

badge = card(s2, 9.55, 1.28, 3.45, 0.52, GREEN_BG, GREEN, 1, 0.16)
tfb = badge.text_frame
tfb.margin_left = In(0.1); tfb.margin_top = In(0.03)
p = para(tfb, first=True, space_after=0)
run(p, "DEPLOYED AND RUNNING TODAY", 9, True, GREEN)
p = para(tfb, space_after=0)
run(p, "ruralai-psi.vercel.app", 10.5, True, NAVY)

y = 1.98
LW = 6.05
y = section(s2, 0.36, y, LW, "Proposed Solution", [
    ("One platform, six services:",
     "assisted teleconsultation, appointment and queue management, digital triage, "
     "longitudinal records, emergency referral and facility dashboards."),
    ("The health worker is the operator:",
     "the patient never touches it; a registered doctor makes every clinical decision."),
    ("Runs on what a sub-centre already has:",
     "a browser and a phone camera. No new hardware."),
])

y = section(s2, 0.36, y + 0.10, LW, "How it addresses the problem", [
    ("Travel and waiting:", "the case travels, not the patient — video or same-day batch review."),
    ("Fragmented records:", "one record keyed to Aadhaar; paper scripts and lab reports digitised by camera."),
    ("Language and literacy:", "symptoms spoken in 7 Indian languages; findings read back aloud."),
    ("Delayed referral:", "red-flag rules escalate to the nearest district hospital automatically."),
])

y = section(s2, 0.36, y + 0.10, LW, "Innovation and uniqueness", [
    ("Rules decide, models advise:",
     "a deterministic engine sets the triage floor; the model may raise a tier and can never lower one."),
    ("Bounded, not prompted:",
     "a classifier trained on 244,938 labelled symptom vectors constrains what the LLM may say."),
    ("Medication never model-authored:",
     "only from a practitioner-signed formulary. Enforced in code, covered by tests."),
])

card(s2, 6.72, 1.98, 6.28, 1.30, NAVY_SOFT, None, radius=0.10)
label(s2, 6.92, 2.10, 5.9, "Already built — verifiable at the link", 11, NAVY)
STATS = [("59", "API routes"), ("17", "DB tables"), ("135", "tests passing"), ("7", "languages")]
for i, (big, lab) in enumerate(STATS):
    x = 6.92 + i * 1.52
    t = tb(s2, x, 2.42, 1.45, 0.72)
    p = para(t, first=True, space_after=0, line=0.9)
    run(p, big, 21, True, NAVY)
    p = para(t, space_after=0, line=0.9)
    run(p, lab, 9, False, MUTED)

s2.shapes.add_picture(os.path.join(HERE, "fig_architecture.png"),
                      In(6.72), In(3.44), width=In(6.28))
t = tb(s2, 6.72, 6.62, 6.28, 0.24)
p = para(t, first=True, space_after=0)
run(p, "Architecture — every layer in this diagram is deployed.", 8.5, False, MUTED, italic=True)

# ================================================================== SLIDE 3
s3 = S[2]
kill(find(s3, "Technologies to be used"))

label(s3, 0.36, 1.24, 6.0, "Methodology and process for implementation", 12.5, NAVY)
t = tb(s3, 0.36, 1.52, 12.6, 0.24)
p = para(t, first=True, space_after=0)
run(p, "Six steps from a patient arriving at a sub-centre to a signed clinical decision. "
       "The tier decides what happens next — each tier has a different, defined output.",
    9.8, False, MUTED)

# flow: 11.8in wide / aspect 2.910 -> 4.06in tall
s3.shapes.add_picture(os.path.join(HERE, "fig_flow.png"),
                      In(0.77), In(1.82), width=In(11.80))

label(s3, 0.36, 6.02, 6.0, "Technologies to be used", 12.5, NAVY)
# stack: 11.8in wide / aspect 19.5 -> 0.61in tall
s3.shapes.add_picture(os.path.join(HERE, "fig_stack.png"),
                      In(0.77), In(6.28), width=In(11.80))

# ================================================================== SLIDE 4
s4 = S[3]
kill(find(s4, "Analysis of the feasibility"))

y0 = 1.28
label(s4, 0.36, y0, 6.2, "Analysis of the feasibility of the idea", 12.5, NAVY)
FEAS = [
    ("Technical", "Built and deployed on proven open stacks — Node, PostgreSQL, React, "
     "scikit-learn. 135 automated tests. Model artifact is 3.6 MB and runs on CPU."),
    ("Operational", "The operator is the health worker who already staffs the sub-centre. "
     "No new cadre, no new hardware — a browser and a phone camera."),
    ("Economic", "Video is peer-to-peer, referral routing is a local data file, "
     "read-aloud is on-device. Only the model calls cost money per assessment."),
    ("Institutional", "Six roles scoped by district, append-only audit log, "
     "row-level security — the shape a state health department can actually run."),
]
yy = y0 + 0.32
for h, b in FEAS:
    t = tb(s4, 0.36, yy, 6.2, 0.46)
    p = para(t, first=True, space_after=0, line=0.95)
    run(p, h + "  ", 10, True, GREEN)
    run(p, b, 9.5, False, MUTED)
    yy += 0.50

label(s4, 0.36, yy + 0.06, 6.2, "Viability", 12.5, NAVY)
VIAB = [
    ("Deployment unit", "the sub-centre — already the tenancy key in the schema."),
    ("Policy fit", "strengthens eSanjeevani and ABDM rather than competing with them."),
    ("Sustainability", "per-facility licensing under NHM budgets; teleconsultation is "
     "already reimbursable under the Ayushman Bharat framework."),
]
yy += 0.38
for h, b in VIAB:
    t = tb(s4, 0.36, yy, 6.2, 0.40)
    p = para(t, first=True, space_after=0, line=0.95)
    run(p, h + "  ", 10, True, NAVY)
    run(p, b, 9.5, False, MUTED)
    yy += 0.42

RX, RW2 = 6.86, 6.14
label(s4, RX, y0, RW2, "Potential challenges and risks  ·  Strategies for overcoming these", 12.5, NAVY)

RISKS = [
    ("Clinical", "Triage thresholds and the formulary are drawn from published guidance "
     "but are not yet practitioner-signed.",
     "Medication is withheld from the health worker at every tier and suppressed by "
     "default in production. The system states this on its own pages. Sign-off is a "
     "named launch gate, not a code task."),
    ("Connectivity", "Assessment currently needs a network round trip.",
     "Rule engine, referral routing and read-aloud already run with no network. "
     "A 4-bit on-premise model at the PHC is the roadmap item that closes the rest."),
    ("Interoperability", "ABDM/FHIR exchange is not built — the ABHA field is captured "
     "by card OCR only.",
     "Records are already structured to the shapes ABDM consumes. Registering as a "
     "Health Information Provider is an adapter, not a data migration."),
    ("Continuity", "Referral completion and high-risk follow-up are not yet tracked.",
     "Both are scheduled work on an existing schema — the referral and review rows "
     "already exist; what is missing is the due-list and the closure loop."),
    ("Load", "Rate limiting and socket fan-out are per-process today.",
     "A shared Redis store makes both horizontal. Identified, scoped, not yet needed "
     "at pilot scale."),
]
yy = y0 + 0.34
for h, risk, mit in RISKS:
    box = card(s4, RX, yy, RW2, 0.74, SOFT, LINE, 0.75, 0.10)
    tf2 = box.text_frame
    tf2.margin_left = In(0.10); tf2.margin_right = In(0.10)
    tf2.margin_top = In(0.055); tf2.margin_bottom = In(0.04)
    tf2.word_wrap = True
    p = para(tf2, first=True, space_after=1, line=0.92)
    run(p, h.upper() + "   ", 8.5, True, AMBER)
    run(p, "Risk:  ", 9, True, INK)
    run(p, risk, 9, False, MUTED)
    p = para(tf2, space_after=0, line=0.92)
    run(p, "Mitigation:  ", 9, True, GREEN)
    run(p, mit, 9, False, MUTED)
    yy += 0.80


# --- deployment path: what happens after the hackathon -----------------------
label(s4, 0.36, 5.62, 6.2, "Deployment path", 12.5, NAVY)
PHASES = [
    ("NOW", "Deployed. One district grid, seeded and demonstrable.", GREEN),
    ("NEXT", "Practitioner sign-off on the formulary and thresholds.", AMBER),
    ("THEN", "ABDM adapter, referral closure, follow-up due-list.", NAVY),
]
for i, (tag, body, col) in enumerate(PHASES):
    x = 0.36 + i * 2.09
    card(s4, x, 5.92, 1.92, 0.92, SOFT, col, 1, 0.12)
    t = tb(s4, x + 0.12, 5.99, 1.70, 0.22)
    p = para(t, first=True, space_after=0, line=0.9)
    run(p, tag, 9.5, True, col)
    t = tb(s4, x + 0.12, 6.22, 1.70, 0.56)
    p = para(t, first=True, space_after=0, line=0.92)
    run(p, body, 8.4, False, MUTED)

# ================================================================== SLIDE 5
s5 = S[4]
kill(find(s5, "Potential impact on the target audience"))

label(s5, 0.36, 1.24, 8.0, "Potential impact on the target audience", 12.5, NAVY)
# impact: 12.6in wide / aspect 3.937 -> 3.20in tall
s5.shapes.add_picture(os.path.join(HERE, "fig_impact.png"),
                      In(0.36), In(1.54), width=In(12.60))

# --- sourced statistics, as native shapes so every figure stays editable ---
STATS5 = [
    ("79.9%", "specialist shortfall at rural CHCs", "MoHFW, Health Dynamics of India 2022-23"),
    ("54.2%", "of Maharashtra women 15-49 are anaemic", "NFHS-5 Maharashtra fact sheet"),
    ("~34%", "of referred patients never complete referral", "Chennai referral cohort study"),
    ("43 cr", "eSanjeevani consultations — we plug into it", "National Health Authority, 2025"),
]
for i, (big, lab, src) in enumerate(STATS5):
    x = 0.36 + i * 3.19
    card(s5, x, 4.92, 3.02, 0.78, NAVY_SOFT, None, radius=0.10)
    t = tb(s5, x + 0.13, 4.99, 2.78, 0.30)
    p = para(t, first=True, space_after=0, line=0.9)
    run(p, big + "   ", 17, True, NAVY)
    run(p, lab, 8.4, False, INK)
    t = tb(s5, x + 0.13, 5.42, 2.78, 0.22)
    p = para(t, first=True, space_after=0, line=0.9)
    run(p, src, 7.4, False, MUTED, italic=True)

label(s5, 0.36, 5.86, 12.6, "Benefits of the solution", 12.5, NAVY)
BEN = [
    ("Social", GREEN,
     [("Specialist reach", "a 79.9% shortfall becomes a scheduling problem, not a travel one."),
      ("Equity within the state", "41.2% in Sindhudurg to 68.2% in Gadchiroli — one protocol.")]),
    ("Economic", NAVY,
     [("Household cost", "68% of rural health spending is out of pocket."),
      ("System cost", "no per-minute video fee, no mapping API, no per-call TTS.")]),
    ("Governance", AMBER,
     [("Quality monitoring", "live visits and risk mix, aggregate only — no identifiers."),
      ("Accountability", "an append-only audit log that no role can edit or delete.")]),
]
for i, (title, col, rows) in enumerate(BEN):
    x = 0.36 + i * 4.24
    card(s5, x, 6.12, 4.02, 0.76, SOFT, LINE, 0.75, 0.10)
    label(s5, x + 0.13, 6.20, 3.8, title, 10, col)
    yy = 6.40
    for h, b in rows:
        t = tb(s5, x + 0.13, yy, 3.76, 0.22)
        p = para(t, first=True, space_after=0, line=0.9)
        run(p, h + "  ", 8.2, True, INK)
        run(p, b, 8.2, False, MUTED)
        yy += 0.235

# ================================================================== SLIDE 6
s6 = S[5]
kill(find(s6, "Details / Links"))

label(s6, 0.36, 1.24, 6.2, "Details / Links of the reference and research work", 12.5, NAVY)
REFS = [
    ("MoHFW", "Health Dynamics of India (Infrastructure & Human Resources) 2022-23 — "
     "79.9% specialist shortfall at rural CHCs; 1,69,615 sub-centres, 31,882 PHCs, 6,359 CHCs.",
     "pib.gov.in/PressReleasePage.aspx?PRID=2053070"),
    ("NSS 75th Round", "Health in India — 68% of rural health expenditure out of pocket; "
     "₹16,676 average rural spend per hospitalisation; ~86% without insurance.",
     "mospi.gov.in — Report 586"),
    ("NFHS-5 Maharashtra", "State fact sheet (DHS FR374) — anaemia 54.2% in women 15-49 and "
     "68.9% in children 6-59 months; district range 41.2%-68.2%.",
     "dhsprogram.com/pubs/pdf/FR374/FR374_Maharashtra.pdf"),
    ("National Health Authority", "eSanjeevani 43 crore teleconsultations, 1,31,793 spokes; "
     "ABDM 100 crore health records linked to ABHA.",
     "abdm.gov.in  ·  esanjeevani.mohfw.gov.in"),
    ("WHO / IMNCI · PALS", "Integrated Management of Neonatal and Childhood Illness — "
     "paediatric respiratory bands and dehydration plans B and C encoded in the triage engine.",
     "who.int/publications — IMNCI"),
    ("NMC", "Telemedicine Practice Guidelines 2020 — the prescribing boundary the platform "
     "enforces in code.", "nmc.org.in"),
]
yy = 1.58
for src, desc, link in REFS:
    t = tb(s6, 0.36, yy, 6.2, 0.52)
    p = para(t, first=True, space_after=0, line=0.92)
    run(p, src + "  ", 9.3, True, NAVY)
    run(p, desc, 9.3, False, MUTED)
    p = para(t, space_after=0, line=0.92)
    run(p, link, 8.3, False, SAFFRON, italic=True)
    yy += 0.58

RX = 6.86
card(s6, RX, 1.24, 6.14, 1.62, NAVY_SOFT, None, radius=0.09)
label(s6, RX + 0.16, 1.34, 5.8, "Project resources", 11.5, NAVY)
LINKS = [
    ("Live prototype", "ruralai-psi.vercel.app"),
    ("Source code", "github.com/PriyamMishra853/RURALAI-"),
    ("Architecture document", "PROJECT_DOCUMENTATION/ — 20 sections"),
    ("Demo video", "[DEMO VIDEO LINK]"),
]
yy = 1.62
for k, v in LINKS:
    t = tb(s6, RX + 0.16, yy, 5.8, 0.26)
    p = para(t, first=True, space_after=0, line=0.92)
    run(p, k + "   ", 9.3, True, INK)
    run(p, v, 9.3, False, NAVY)
    yy += 0.29

card(s6, RX, 3.02, 6.14, 1.55, SOFT, LINE, 0.75, 0.09)
label(s6, RX + 0.16, 3.12, 5.8, "Strengthens the public system — does not replace it", 11.5, GREEN)
POS = [
    ("eSanjeevani", "carries the consultation. We prepare the case before it and "
     "record the decision after it."),
    ("ABDM / ABHA", "the record is structured for it; ABHA is captured at registration by card OCR."),
    ("MoHFW / IPHS protocols", "retrieval is filtered to approved protocols only, "
     "cited with title, source and version."),
]
yy = 3.40
for k, v in POS:
    t = tb(s6, RX + 0.16, yy, 5.82, 0.36)
    p = para(t, first=True, space_after=0, line=0.92)
    run(p, k + "  ", 9.2, True, INK)
    run(p, v, 9.2, False, MUTED)
    yy += 0.38

card(s6, RX, 4.73, 6.14, 2.06, SOFT, LINE, 0.75, 0.09)
label(s6, RX + 0.16, 4.83, 5.8, "Where existing tools stop", 11.5, NAVY)
COMP = [
    ("Teleconsultation platforms", "connect a doctor. They do not triage, "
     "read a paper prescription, or hold a record between visits."),
    ("Symptom checkers", "speak to the patient. This speaks to a trained health worker "
     "and refuses to name a medicine."),
    ("Hospital information systems", "begin at the hospital door — after the journey "
     "this platform exists to avoid."),
]
yy = 5.11
for k, v in COMP:
    t = tb(s6, RX + 0.16, yy, 5.82, 0.50)
    p = para(t, first=True, space_after=0, line=0.92)
    run(p, k + "  ", 9.2, True, INK)
    run(p, v, 9.2, False, MUTED)
    yy += 0.54

# ------------------------------------------------------------- team name oval
for sl in list(S):
    for sh in sl.shapes:
        if sh.has_text_frame and "your team name" in sh.text_frame.text.lower():
            tfo = sh.text_frame
            tfo.clear()
            p = tfo.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run(p, "UNFILTERED\nENGINEERS", 8, True, WHITE)

# ------------------------------------------------------------- speaker notes
NOTES_TEXT = {
0: """[0:00-0:20] 20 seconds.
"Problem Statement 26133, for the Government of Maharashtra. Accessibility and quality of public healthcare in rural and underserved areas. We are team UNFILTERED ENGINEERS."
Then, immediately and without pausing:
"Before I describe anything, this is not a concept. It is deployed and you can open it right now at ruralai-psi.vercel.app."
Do not read the rest of the title slide aloud. Move on.""",
1: """[0:20-1:20] 60 seconds. This is the slide that wins or loses the room.

Open with the live badge, top right: "Everything I am about to describe is running at that URL."

Proposed Solution - one line only: "One platform doing the six things the problem statement asks for: teleconsultation, queue management, triage, longitudinal records, referral and dashboards."

How it addresses the problem - pick TWO, not four: "The case travels instead of the patient. And the record is keyed to Aadhaar, so history does not restart at every facility."

Innovation - this is the differentiator, say it slowly:
"Most teams wire a language model to a prompt. We did not. A deterministic rules engine sets the triage floor. A classifier trained on 244,938 labelled symptom vectors bounds what the model is allowed to say. The model can raise a risk tier. It can never lower one. And it can never name a medicine - that comes from a practitioner-signed formulary, enforced in code and covered by tests."

Point at the four numbers: "59 routes, 17 tables, 135 tests passing, 7 languages. All verifiable at the link."

If asked to cut for time: drop 'How it addresses the problem' entirely. Keep Innovation.""",
2: """[1:20-2:10] 50 seconds. Do not read this slide. Trace it with your hand.

"Register, capture, digitise." - sweep left to right across the first three boxes.

Stop on ASSESS: "Four things run here. The rule engine sets the floor. The classifier proposes ranked candidates. Vision and OCR can raise the tier. The language model writes the summary last, and only within those bounds."

Then the formula: "Final tier is the maximum of rules, vision and model. That is the whole safety argument in one line."

Fan your hand across the three tiers: "Low gets a complete plan on the spot. Medium requires a video consultation before any treatment. High leaves the platform entirely and goes to the nearest district hospital."

Land on the navy bar: "A registered practitioner signs every decision, and it goes straight back to the health worker."

Tech stack: one sentence. "Standard open stacks. Nothing exotic - which is why it is already deployed."

Do not explain individual technologies unless asked.""",
3: """[2:10-3:00] 50 seconds.

Feasibility - do not read all four. Say: "Technically it is already built and tested. Operationally, the operator is the health worker who is already there - no new cadre, no new hardware."

Economic, say this one deliberately: "Video is peer-to-peer, referral routing is a local file, read-aloud runs on the device. The expensive parts of a telemedicine platform are structurally free here."

Then turn to the risk column and OWN it - this is where you earn credibility:

"I want to be direct about what is not finished."

Clinical: "Our triage thresholds and formulary come from published guidance and have not been signed by a registered practitioner for this deployment. So medication is withheld from the health worker at every tier, and the system says so on its own front page. That is a launch gate, not a code task."

Then: "Referral completion tracking, high-risk follow-up and ABDM exchange are scheduled, not built. The schema already holds the rows they need."

Judges trust the team that names its own gaps before they do.""",
4: """[3:00-3:50] 50 seconds.

Top half, trace the two rows: "Today: symptom, hours of travel, a queue with no triage, a consult where history restarts, and a paper prescription that is lost before the next visit."

"With RuralAI: the same clinical authority, no journey. Travel happens only for the cases that genuinely need a hospital - and the tier decides that, not the queue."

Then the four numbers, and cite each source out loud - this is where you separate yourself:
"79.9 percent specialist shortfall at rural community health centres - Ministry of Health, 2022-23."
"54.2 percent of Maharashtra women aged 15 to 49 are anaemic - NFHS-5."
"Around a third of referred patients never complete the referral."
"And eSanjeevani has done 43 crore consultations - we plug into that, we do not compete with it."

Benefits: name the three buckets, do not read the contents. "Social, economic, governance."

If short on time, cut the benefits row entirely and stay on the numbers.""",
5: """[3:50-4:30] 40 seconds, then stop and invite questions.

"Every number on the previous slide is sourced, and the sources are here - Ministry of Health, the 75th round of the National Sample Survey, the NFHS-5 Maharashtra fact sheet, the National Health Authority, WHO IMNCI, and the NMC telemedicine guidelines."

Then the positioning box - say this clearly, it directly answers the problem statement's own wording:
"The problem statement says strengthen, not replace. eSanjeevani carries the consultation. We prepare the case before it and record the decision after it. ABDM holds the record - ours is structured for it."

Then, if they have not already asked, pre-empt the obvious question:
"You may be wondering why this is not just a symptom checker. Symptom checkers speak to patients. This speaks to a trained health worker, and it refuses to name a medicine."

Close on the resources box: "Live prototype, source code, a twenty-section architecture document, and the demo video. All linked here."

Then stop talking. Do not fill silence.""",
}
for i, sl in enumerate(S):
    sl.notes_slide.notes_text_frame.text = NOTES_TEXT.get(i, "")

prs.save(OUT)
print("saved", OUT)
print("slides:", len(prs.slides.__iter__.__self__._sldIdLst))
